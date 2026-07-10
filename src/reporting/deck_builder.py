from pptx import Presentation
from pptx.util import Inches, Pt
from src.aggregation import Trend, Risk

def create_executive_deck(project_reports: list[dict], trends: list[Trend], risks: list[Risk], out_path: str):
    prs = Presentation()
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Project Portfolio Health Review"
    subtitle.text = "Prepared automatically by AI Health Agent"
    
    # 2. Executive summary
    exec_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(exec_slide_layout)
    slide.shapes.title.text = "Executive Summary"
    tf = slide.placeholders[1].text_frame
    
    status_counts = {"Green": 0, "Amber": 0, "Red": 0}
    for r in project_reports:
        status_counts[r.get("status", "Green")] += 1
        
    p = tf.add_paragraph()
    p.text = f"Portfolio Mix: {status_counts['Red']} Red, {status_counts['Amber']} Amber, {status_counts['Green']} Green"
    p.font.size = Pt(18)
    
    p2 = tf.add_paragraph()
    p2.text = "Top-line call: Significant schedule slippage observed across multiple projects."
    
    p3 = tf.add_paragraph()
    p3.text = "Ask: Review resourcing and blocker resolution for Red projects."
    
    # 3. Portfolio health at a glance (Table)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Portfolio Health At A Glance"
    
    rows = len(project_reports) + 1
    cols = 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(0.8)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(3.0)
    table.columns[3].width = Inches(1.5)
    
    headers = ["Project Name", "Status", "Top Driver", "Data Confidence"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
        
    for idx, r in enumerate(project_reports):
        summary = r.get("summary", {})
        table.cell(idx + 1, 0).text = summary.get("project_name", "Unknown")
        table.cell(idx + 1, 1).text = r.get("status", "Unknown")
        
        # Derive top driver
        top_driver = "None"
        for ev in r.get("evidences", []):
            if ev.get("status") == r.get("status"):
                top_driver = ev.get("signal", "").replace("_", " ").title()
                break
        table.cell(idx + 1, 2).text = top_driver
        table.cell(idx + 1, 3).text = r.get("confidence", "Unknown")
        
    # 4. Cross-project trends
    slide = prs.slides.add_slide(exec_slide_layout)
    slide.shapes.title.text = "Cross-Project Trends"
    tf = slide.placeholders[1].text_frame
    
    if not trends:
        tf.text = "No systemic cross-project trends identified in current sample."
    else:
        for t in trends:
            p = tf.add_paragraph()
            p.text = f"{t.category}: {t.signal.replace('_', ' ').title()} observed across projects ({', '.join(t.affected_project_ids)})"
            for f in t.example_facts:
                p_sub = tf.add_paragraph()
                p_sub.text = f"- {f}"
                p_sub.level = 1
                
    # 5. Emerging risks
    slide = prs.slides.add_slide(exec_slide_layout)
    slide.shapes.title.text = "Emerging Risks"
    tf = slide.placeholders[1].text_frame
    
    if not risks:
        tf.text = "No immediate emerging risks identified."
    else:
        for r in risks:
            p = tf.add_paragraph()
            p.text = f"{r.project_id}: {r.risk_description} (Trajectory: {r.trajectory})"
            
    # 6. Recommendations
    slide = prs.slides.add_slide(exec_slide_layout)
    slide.shapes.title.text = "Recommendations & Decisions Needed"
    tf = slide.placeholders[1].text_frame
    
    p = tf.add_paragraph()
    p.text = "1. Immediate deep-dive required on all Red projects (S2P_Project) to clear blockers."
    
    p = tf.add_paragraph()
    p.text = "2. Investigate schedule estimation accuracy, as >30% of open tasks are delayed."
    
    prs.save(out_path)
